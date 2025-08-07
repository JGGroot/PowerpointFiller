import streamlit as st
import pandas as pd
from pptx import Presentation
from pptx.util import Inches
from PIL import Image
import json
import re
import io
import zipfile
from datetime import datetime
import base64
from clipboard_component import copy_component, paste_component
import docx
import glob
import os

# PDF support imports
import fitz  # PyMuPDF
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
import PyPDF2

# Configure the page
st.set_page_config(
    page_title="Document AI Field Filler",
    page_icon="📊",
    layout="centered",
    initial_sidebar_state="collapsed"
)

# Load prompt configuration
@st.cache_data
def load_prompt_config():
    """Load prompt configuration from config file"""
    config_file = "prompt_config.json"
    default_config = {
        "default_prompt": """I need you to analyze project data and extract information for specific document fields. Return ONLY a valid JSON object with the field names as keys and extracted values as values.
**Document Fields to Fill:**

{field_list}

**Instructions:**

1. Extract relevant information from the data for each field
2. If a field name suggests specific content (e.g., "commander_name" should be a person's name), extract accordingly
3. Be clear, professional, and concise. You are drafting documents for official government use so no slang etc. 
4. Conduct market research with a focus on Department of Defense, Department of the Air Force, and with the goals of the 100th ARW and 352nd SOW mission goals in mind
5. For fields with money, phone numbers, or other implied formatting, format the extracted values accordingly
6. For fields you can't determine from the data, use "TBD" or leave reasonable placeholder text based on context
7. Return ONLY the JSON object - no explanations or additional text

**Project Data to Analyze:**

{project_data}

Please analyze the above data and return the JSON object with field values""",
        "template_prompts": {
            "example_template.pptx": "Custom prompt for example template...",
            "military_brief.pptx": "Military-specific prompt for briefing template..."
        }
    }
    
    try:
        if os.path.exists(config_file):
            with open(config_file, 'r', encoding='utf-8') as f:
                config = json.load(f)
                # Ensure default_prompt exists
                if "default_prompt" not in config:
                    config["default_prompt"] = default_config["default_prompt"]
                # Ensure template_prompts exists
                if "template_prompts" not in config:
                    config["template_prompts"] = {}
                return config
        else:
            # Create default config file
            with open(config_file, 'w', encoding='utf-8') as f:
                json.dump(default_config, f, indent=2)
            # st.info(f"Created default prompt configuration file: {config_file}")
            return default_config
    except Exception as e:
        st.warning(f"Could not load prompt config: {e}. Using default prompt.")
        return default_config

def get_template_prompt(template_name, prompt_config):
    """Get the appropriate prompt for a template"""
    if template_name and template_name in prompt_config.get("template_prompts", {}):
        st.info(f"Using custom prompt for template: {template_name}")
        return prompt_config["template_prompts"][template_name]
    else:
        if template_name and template_name != "Upload my own template":
            st.info(f"No custom prompt found for '{template_name}', using default prompt")
        return prompt_config["default_prompt"]

# Custom CSS for better styling
st.markdown("""
<style>
    .main-header {
        text-align: center;
        color: #2c3e50;
        font-size: 2.5rem;
        margin-bottom: 2rem;
        padding: 1rem;
        background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
        color: white;
        border-radius: 10px;
    }
    .step-container {
        background: #330066;
        padding: 20px;
        border-radius: 10px;
        margin: 20px 0;
        border-left: 5px solid #007bff;
    }
    .success-box {
        background-color: #d4edda;
        border: 1px solid #c3e6cb;
        border-radius: 5px;
        padding: 15px;
        margin: 10px 0;
    }
    .warning-box {
        background-color: #fff000;
        border: 1px solid #ffeaa7;
        border-radius: 5px;
        padding: 15px;
        margin: 10px 0;
    }
    .ai-button {
        margin: 5px;
        padding: 10px 20px;
        border-radius: 5px;
        border: none;
        color: white;
        font-weight: bold;
        text-decoration: none;
        display: inline-block;
        text-align: center;
    }
    .nipr-btn { background-color: #004d40; } /* Dark Teal for NiprGPT */
</style>
""", unsafe_allow_html=True)

# Google Analytics - Replace 'G-HMVVJJ6C17' with your actual Google Analytics ID
st.markdown("""
<!-- Google tag (gtag.js) -->
<script async src="https://www.googletagmanager.com/gtag/js?id=G-HMVVJJ6C17"></script>
<script>
  window.dataLayer = window.dataLayer || [];
  function gtag(){dataLayer.push(arguments);}
  gtag('js', new Date());
  gtag('config', 'G-HMVVJJ6C17');
</script>
""", unsafe_allow_html=True)

# --- Analysis Functions ---

def analyze_pdf_fields(uploaded_file):
    """Analyze PDF file for field placeholders and form fields"""
    try:
        # Reset file pointer
        if hasattr(uploaded_file, 'seek'):
            uploaded_file.seek(0)
        
        # Read PDF with PyMuPDF
        pdf_bytes = uploaded_file.read()
        pdf_document = fitz.open(stream=pdf_bytes, filetype="pdf")
        
        found_fields = set()
        field_locations = []
        field_pattern = r'\{\{([^}]+)\}\}'
        
        # Check if PDF is encrypted/protected
        if pdf_document.needs_pass:
            st.warning("⚠️ PDF is password protected. Please provide the password or use an unprotected PDF.")
            pdf_password = st.text_input("Enter PDF password:", type="password", key="pdf_password")
            if pdf_password:
                if pdf_document.authenticate(pdf_password):
                    st.success("✅ PDF unlocked successfully!")
                else:
                    st.error("❌ Incorrect password. Please try again.")
                    pdf_document.close()
                    return [], []
            else:
                pdf_document.close()
                return [], []
        
        # Method 1: Extract text and look for {{field_name}} patterns
        for page_num in range(len(pdf_document)):
            page = pdf_document.load_page(page_num)
            text_content = page.get_text()
            
            # Find field patterns in text
            matches = re.findall(field_pattern, text_content)
            for field in matches:
                found_fields.add(field)
                field_locations.append({
                    'field': field,
                    'page': page_num + 1,
                    'type': 'text',
                    'context': text_content[:100] + '...' if len(text_content) > 100 else text_content
                })
        
        # Method 2: Check for form fields (if it's a fillable PDF)
        try:
            # Reset file pointer for PyPDF2
            if hasattr(uploaded_file, 'seek'):
                uploaded_file.seek(0)
            
            pdf_reader = PyPDF2.PdfReader(uploaded_file)
            
            # Handle encrypted PDFs
            if pdf_reader.is_encrypted:
                # Try to decrypt with empty password first (common case)
                try:
                    pdf_reader.decrypt("")
                    st.info("📋 PDF encryption bypassed for form field analysis.")
                except:
                    # If that fails and we have a password from earlier, try it
                    if 'pdf_password' in st.session_state and st.session_state.pdf_password:
                        try:
                            pdf_reader.decrypt(st.session_state.pdf_password)
                            st.info("📋 Using provided password for form field analysis.")
                        except:
                            st.warning("⚠️ Could not decrypt PDF for form field analysis. Text analysis will still work.")
                            pdf_reader = None
                    else:
                        st.warning("⚠️ PDF is encrypted. Form field detection limited. Text pattern detection will still work.")
                        pdf_reader = None
            
            if pdf_reader:
                # Check each page for form fields
                for page_num, page in enumerate(pdf_reader.pages):
                    if '/Annots' in page:
                        annotations = page['/Annots']
                        if annotations:
                            for annotation_ref in annotations:
                                try:
                                    annotation = annotation_ref.get_object()
                                    if annotation.get('/Subtype') == '/Widget':
                                        field_name = annotation.get('/T')
                                        if field_name:
                                            field_name_str = str(field_name)
                                            # Check if field name contains our pattern
                                            pattern_matches = re.findall(field_pattern, field_name_str)
                                            if pattern_matches:
                                                for field in pattern_matches:
                                                    found_fields.add(field)
                                                    field_locations.append({
                                                        'field': field,
                                                        'page': page_num + 1,
                                                        'type': 'form_field_pattern',
                                                        'field_name': field_name_str
                                                    })
                                            else:
                                                # Add the form field name itself as a potential field
                                                found_fields.add(field_name_str)
                                                field_locations.append({
                                                    'field': field_name_str,
                                                    'page': page_num + 1,
                                                    'type': 'form_field',
                                                    'field_name': field_name_str
                                                })
                                except Exception as field_error:
                                    # Skip problematic form fields
                                    continue
                                    
        except Exception as e:
            st.warning(f"Form field analysis encountered issues: {e}. Text analysis completed successfully.")
        
        pdf_document.close()
        return list(found_fields), field_locations
        
    except Exception as e:
        st.error(f"Error analyzing PDF: {str(e)}")
        st.info("💡 **PDF Troubleshooting Tips:**\n- Ensure the PDF is not corrupted\n- Try removing password protection\n- Check if text is selectable (not scanned image)")
        return [], []

def analyze_powerpoint_fields(uploaded_file):
    """(Corrected) Analyze PowerPoint file for field placeholders"""
    try:
        prs = Presentation(uploaded_file)
        found_fields = set()
        field_locations = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            for shape in slide.shapes:
                if hasattr(shape, "text_frame") and shape.text_frame and shape.text_frame.text:
                    text_content = shape.text_frame.text
                    field_pattern = r'\{\{([^}]+)\}\}'
                    matches = re.findall(field_pattern, text_content)
                    
                    for field in matches:
                        found_fields.add(field)
                        field_locations.append({
                            'field': field,
                            'slide': slide_num,
                            'context': text_content[:100] + '...' if len(text_content) > 100 else text_content
                        })
                
                elif shape.has_table:
                    table = shape.table
                    for row_num, row in enumerate(table.rows):
                        for cell_num, cell in enumerate(row.cells):
                            if cell.text:
                                text_content = cell.text
                                field_pattern = r'\{\{([^}]+)\}\}'
                                matches = re.findall(field_pattern, text_content)
                                
                                for field in matches:
                                    found_fields.add(field)
                                    field_locations.append({
                                        'field': field,
                                        'slide': slide_num,
                                        'location': f'Table R{row_num+1}C{cell_num+1}',
                                        'context': text_content[:50] + '...' if len(text_content) > 50 else text_content
                                    })
        
        return list(found_fields), field_locations

    except Exception as e:
        st.error(f"Error analyzing PowerPoint: {str(e)}")
        return [], []

def analyze_word_fields(uploaded_file):
    """(FIXED) Analyze a Word document for field placeholders including text boxes."""
    try:
        doc = docx.Document(uploaded_file)
        found_fields = set()
        field_pattern = r'\{\{([^}]+)\}\}'

        # Check regular paragraphs
        for para in doc.paragraphs:
            if para.text:
                matches = re.findall(field_pattern, para.text)
                for field in matches:
                    found_fields.add(field)

        # Check tables
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    for para in cell.paragraphs:
                        if para.text:
                            matches = re.findall(field_pattern, para.text)
                            for field in matches:
                                found_fields.add(field)

        # Check text boxes and shapes (NEW CODE)
        try:
            from docx.oxml.ns import nsdecls, qn
            from lxml import etree
            
            # Define namespaces explicitly
            namespaces = {
                'w': 'http://schemas.openxmlformats.org/wordprocessingml/2006/main',
                'wp': 'http://schemas.openxmlformats.org/drawingml/2006/wordprocessingDrawing',
                'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'
            }
            
            # Get the document's XML tree
            doc_xml = doc.element
            
            # Look for text boxes using various possible paths
            textbox_paths = [
                './/w:drawing//w:txbxContent//w:p//w:t',
                './/w:drawing//a:txBody//a:p//a:t', 
                './/w:object//w:drawing//w:txbxContent//w:p//w:t',
                './/w:pict//w:textbox//w:txbxContent//w:p//w:t'
            ]
            
            for path in textbox_paths:
                try:
                    text_elements = doc_xml.xpath(path, namespaces=namespaces)
                    for t_elem in text_elements:
                        if t_elem.text:
                            matches = re.findall(field_pattern, t_elem.text)
                            for field in matches:
                                found_fields.add(field)
                except Exception:
                    # Skip this path if it doesn't work
                    continue
            
            # Alternative approach: search through all text elements in the document
            # This is a broader search that should catch text boxes regardless of structure
            try:
                all_text_elements = doc_xml.xpath('.//w:t', namespaces=namespaces)
                for t_elem in all_text_elements:
                    if t_elem.text:
                        matches = re.findall(field_pattern, t_elem.text)
                        for field in matches:
                            found_fields.add(field)
            except Exception:
                pass
                
        except Exception as e:
            # If XML parsing fails, fall back to a simpler approach
            st.warning(f"Advanced text box detection failed: {e}. Using basic detection only.")

        # Alternative approach: check headers and footers which might contain text boxes
        for section in doc.sections:
            # Check headers
            if section.header:
                for para in section.header.paragraphs:
                    if para.text:
                        matches = re.findall(field_pattern, para.text)
                        for field in matches:
                            found_fields.add(field)
            
            # Check footers  
            if section.footer:
                for para in section.footer.paragraphs:
                    if para.text:
                        matches = re.findall(field_pattern, para.text)
                        for field in matches:
                            found_fields.add(field)
        
        return list(found_fields), []
        
    except Exception as e:
        st.error(f"Error analyzing Word document: {e}")
        return [], []

# --- Helper and Filling Functions ---

def generate_ai_prompt(fields, project_data, template_name=None):
    """Generate AI prompt using template-specific or default prompt"""
    # Load prompt configuration
    prompt_config = load_prompt_config()
    
    # Get the appropriate prompt template
    prompt_template = get_template_prompt(template_name, prompt_config)
    
    # Prepare field list
    field_descriptions = [f"  - {field}" for field in sorted(fields)]
    field_list = chr(10).join(field_descriptions)
    
    # Replace placeholders in the prompt template
    final_prompt = prompt_template.format(
        field_list=field_list,
        project_data=project_data
    )
    
    return final_prompt

def replace_text_in_paragraph(paragraph, key, value):
    """(NEW) Replaces text in a paragraph, preserving formatting.
    This is a robust function for both pptx and docx.
    """
    if key not in paragraph.text:
        return
    
    # Replace the simple case
    for run in paragraph.runs:
        if key in run.text:
            run.text = run.text.replace(key, value)
            return

    # Handle cases where the key is split across multiple runs
    runs = paragraph.runs
    full_text = "".join(run.text for run in runs)
    if key in full_text:
        start_index = full_text.find(key)
        end_index = start_index + len(key)
        
        current_pos = 0
        runs_to_modify = []
        for run in runs:
            run_len = len(run.text)
            if current_pos < end_index and current_pos + run_len > start_index:
                runs_to_modify.append(run)
            current_pos += run_len
        
        if runs_to_modify:
            # Replace text in the first run and clear others
            original_text = "".join(run.text for run in runs_to_modify)
            new_text = original_text.replace(key, value, 1)
            runs_to_modify[0].text = new_text
            for i in range(1, len(runs_to_modify)):
                runs_to_modify[i].text = ""

def fill_pdf_with_data(pdf_file, data):
    """Fill PDF with data - prioritizing form field filling over text replacement"""
    try:
        # Reset file pointer
        if hasattr(pdf_file, 'seek'):
            pdf_file.seek(0)
        
        replacements_made = 0
        
        # Method 1: Try form field filling first (this is the proper way for Acrobat forms)
        try:
            pdf_reader = PyPDF2.PdfReader(pdf_file)
            pdf_writer = PyPDF2.PdfWriter()
            
            # Handle encryption for form filling
            if pdf_reader.is_encrypted:
                try:
                    # Try empty password first
                    pdf_reader.decrypt("")
                except:
                    # Try with provided password
                    if 'pdf_password' in st.session_state and st.session_state.pdf_password:
                        try:
                            pdf_reader.decrypt(st.session_state.pdf_password)
                        except:
                            st.warning("Could not decrypt PDF for form field filling.")
                            pdf_reader = None
                    else:
                        pdf_reader = None
            
            form_fields_filled = 0
            all_form_fields = set()
            
            if pdf_reader:
                # First pass: collect all form field names
                for page in pdf_reader.pages:
                    if '/Annots' in page:
                        annotations = page['/Annots']
                        if annotations:
                            for annotation_ref in annotations:
                                try:
                                    annotation = annotation_ref.get_object()
                                    if annotation.get('/Subtype') == '/Widget':
                                        field_name = annotation.get('/T')
                                        if field_name:
                                            all_form_fields.add(str(field_name))
                                except:
                                    continue
                
                st.info(f"Found {len(all_form_fields)} form fields in PDF: {list(all_form_fields)}")
                
                # Second pass: fill form fields
                for page in pdf_reader.pages:
                    if '/Annots' in page:
                        annotations = page['/Annots']
                        if annotations:
                            for annotation_ref in annotations:
                                try:
                                    annotation = annotation_ref.get_object()
                                    if annotation.get('/Subtype') == '/Widget':
                                        field_name = annotation.get('/T')
                                        if field_name:
                                            field_name_str = str(field_name)
                                            
                                            # Check if we have data for this field
                                            field_filled = False
                                            for field, value in data.items():
                                                # Direct match
                                                if field == field_name_str:
                                                    annotation.update({
                                                        PyPDF2.generic.NameObject('/V'): 
                                                        PyPDF2.generic.TextStringObject(str(value))
                                                    })
                                                    form_fields_filled += 1
                                                    field_filled = True
                                                    st.success(f"✅ Filled form field '{field_name_str}' with '{value}'")
                                                    break
                                                
                                                # Pattern match {{field_name}}
                                                placeholder = f"{{{{{field}}}}}"
                                                if placeholder in field_name_str:
                                                    annotation.update({
                                                        PyPDF2.generic.NameObject('/V'): 
                                                        PyPDF2.generic.TextStringObject(str(value))
                                                    })
                                                    form_fields_filled += 1
                                                    field_filled = True
                                                    st.success(f"✅ Filled form field '{field_name_str}' with '{value}' (pattern match)")
                                                    break
                                            
                                            if not field_filled:
                                                st.info(f"ℹ️ Form field '{field_name_str}' - no matching data found")
                                                
                                except Exception as form_error:
                                    st.warning(f"Could not process form field: {form_error}")
                                    continue
                    
                    pdf_writer.add_page(page)
                
                if form_fields_filled > 0:
                    # Form fields were filled, now also remove any {{placeholder}} text that might be visible
                    form_output = io.BytesIO()
                    pdf_writer.write(form_output)
                    form_output.seek(0)
                    
                    # Now process with PyMuPDF to remove placeholder text
                    pdf_document = fitz.open(stream=form_output.getvalue(), filetype="pdf")
                    
                    # Handle password-protected PDFs for text removal
                    if pdf_document.needs_pass:
                        if 'pdf_password' in st.session_state and st.session_state.pdf_password:
                            if not pdf_document.authenticate(st.session_state.pdf_password):
                                st.warning("Could not authenticate PDF for placeholder text removal")
                            else:
                                # Remove placeholder text that might still be visible
                                text_removed = 0
                                field_pattern = r'\{\{([^}]+)\}\}'
                                
                                for page_num in range(len(pdf_document)):
                                    page = pdf_document.load_page(page_num)
                                    text_instances = page.get_text("dict")
                                    
                                    for block in text_instances["blocks"]:
                                        if "lines" in block:
                                            for line in block["lines"]:
                                                for span in line["spans"]:
                                                    text = span.get("text", "")
                                                    
                                                    # Check if this text contains placeholder patterns
                                                    if re.search(field_pattern, text):
                                                        try:
                                                            # Remove the placeholder text by covering with white
                                                            rect = fitz.Rect(span["bbox"])
                                                            page.draw_rect(rect, color=(1, 1, 1), fill=(1, 1, 1))
                                                            text_removed += 1
                                                        except Exception as remove_error:
                                                            st.warning(f"Could not remove placeholder text: {remove_error}")
                                
                                if text_removed > 0:
                                    st.info(f"Removed {text_removed} placeholder text instances")
                    
                    # Save the final result
                    final_output = io.BytesIO()
                    pdf_document.save(final_output)
                    pdf_document.close()
                    
                    replacements_made = form_fields_filled
                    st.success(f"✅ Successfully filled {form_fields_filled} form fields and cleaned up placeholders!")
                    return final_output.getvalue(), replacements_made
                else:
                    st.warning("⚠️ No form fields matched your data. Will try text-based replacement as fallback.")
        
        except Exception as e:
            st.warning(f"Form field filling failed: {e}. Trying text replacement fallback.")
        
        # Method 2: Fallback to text replacement (only if form filling failed)
        try:
            if hasattr(pdf_file, 'seek'):
                pdf_file.seek(0)
            
            pdf_bytes = pdf_file.read()
            pdf_document = fitz.open(stream=pdf_bytes, filetype="pdf")
            
            # Handle password-protected PDFs
            if pdf_document.needs_pass:
                if 'pdf_password' in st.session_state and st.session_state.pdf_password:
                    if not pdf_document.authenticate(st.session_state.pdf_password):
                        st.error("Cannot fill PDF: Authentication failed")
                        pdf_document.close()
                        return None, 0
                else:
                    st.error("Cannot fill encrypted PDF without password")
                    pdf_document.close()
                    return None, 0
            
            field_pattern = r'\{\{([^}]+)\}\}'
            text_replacements = 0
            
            # Only do text replacement if no form fields were found
            for page_num in range(len(pdf_document)):
                page = pdf_document.load_page(page_num)
                
                # Get all text instances with their positions
                text_instances = page.get_text("dict")
                
                # Look for field patterns in text
                text_found = False
                for block in text_instances["blocks"]:
                    if "lines" in block:
                        for line in block["lines"]:
                            for span in line["spans"]:
                                text = span.get("text", "")
                                for field, value in data.items():
                                    placeholder = f"{{{{{field}}}}}"
                                    if placeholder in text:
                                        text_found = True
                                        break
                
                if text_found:
                    st.info(f"Found {{field}} patterns in text on page {page_num + 1}. Using text replacement.")
                    
                    # Process text replacements with careful positioning
                    for block in text_instances["blocks"]:
                        if "lines" in block:
                            for line in block["lines"]:
                                for span in line["spans"]:
                                    text = span.get("text", "")
                                    original_text = text
                                    modified_text = text
                                    
                                    for field, value in data.items():
                                        placeholder = f"{{{{{field}}}}}"
                                        if placeholder in modified_text:
                                            modified_text = modified_text.replace(placeholder, str(value))
                                    
                                    if modified_text != original_text:
                                        try:
                                            rect = fitz.Rect(span["bbox"])
                                            font_size = span.get("size", 12)
                                            
                                            # Remove original text first
                                            page.draw_rect(rect, color=(1, 1, 1), fill=(1, 1, 1))
                                            
                                            # Add replacement text
                                            page.insert_text(
                                                rect.top_left,
                                                modified_text,
                                                fontsize=font_size,
                                                color=(0, 0, 0),
                                                fontname="helv"  # Use safe font
                                            )
                                            text_replacements += 1
                                            
                                        except Exception as text_error:
                                            st.warning(f"Could not replace text: {text_error}")
            
            if text_replacements > 0:
                output_buffer = io.BytesIO()
                pdf_document.save(output_buffer)
                pdf_document.close()
                st.success(f"✅ Made {text_replacements} text replacements (fallback method)")
                return output_buffer.getvalue(), text_replacements
            else:
                pdf_document.close()
                st.warning("⚠️ No field patterns found in PDF text or form fields.")
                
        except Exception as e:
            st.error(f"Text replacement also failed: {str(e)}")
        
        return None, 0
        
    except Exception as e:
        st.error(f"Error filling PDF: {str(e)}")
        return None, 0

def fill_powerpoint_with_data(prs, json_data, uploaded_image, progress_container):
    """(CORRECTED) Fill PowerPoint with data preserving formatting."""
    replacements_made = 0
    # Image replacement functionality temporarily disabled
    # if uploaded_image:
    #     # Placeholder for image replacement logic
    #     pass

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape.text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for field, value in json_data.items():
                        placeholder = f"{{{{{field}}}}}"
                        replace_text_in_paragraph(paragraph, placeholder, str(value))
                        replacements_made += 1 # Count attempt
            elif shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for paragraph in cell.text_frame.paragraphs:
                            for field, value in json_data.items():
                                placeholder = f"{{{{{field}}}}}"
                                replace_text_in_paragraph(paragraph, placeholder, str(value))
                                replacements_made += 1 # Count attempt
    return prs, replacements_made

def fill_word_with_data(doc_file, data):
    """(FIXED) Fill a Word document with data, preserving formatting and handling text boxes."""
    doc = docx.Document(doc_file)
    
    # Fill regular paragraphs
    for paragraph in doc.paragraphs:
        for field, value in data.items():
            placeholder = f"{{{{{field}}}}}"
            replace_text_in_paragraph(paragraph, placeholder, str(value))
    
    # Fill tables
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for paragraph in cell.paragraphs:
                    for field, value in data.items():
                        placeholder = f"{{{{{field}}}}}"
                        replace_text_in_paragraph(paragraph, placeholder, str(value))
    
    # Fill text boxes (NEW CODE)
    try:
        from lxml import etree
        
        # Define namespaces explicitly
        namespaces = {
            'w': 'http://schemas.openxmlformats.org/wordprocessingml/2006/main',
            'wp': 'http://schemas.openxmlformats.org/drawingml/2006/wordprocessingDrawing',
            'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'
        }
        
        # Get the document's XML tree
        doc_xml = doc.element
        
        # Look for text elements in text boxes using various paths
        textbox_paths = [
            './/w:drawing//w:txbxContent//w:p//w:t',
            './/w:drawing//a:txBody//a:p//a:t',
            './/w:object//w:drawing//w:txbxContent//w:p//w:t',
            './/w:pict//w:textbox//w:txbxContent//w:p//w:t'
        ]
        
        for path in textbox_paths:
            try:
                text_elements = doc_xml.xpath(path, namespaces=namespaces)
                
                for t_elem in text_elements:
                    if t_elem.text:
                        original_text = t_elem.text
                        modified_text = original_text
                        
                        # Replace all placeholders
                        for field, value in data.items():
                            placeholder = f"{{{{{field}}}}}"
                            modified_text = modified_text.replace(placeholder, str(value))
                        
                        # Update the text if it was modified
                        if modified_text != original_text:
                            t_elem.text = modified_text
                            
            except Exception:
                # Skip this path if it doesn't work
                continue
                
        # Fallback: try to replace in all text elements
        try:
            all_text_elements = doc_xml.xpath('.//w:t', namespaces=namespaces)
            for t_elem in all_text_elements:
                if t_elem.text:
                    original_text = t_elem.text
                    modified_text = original_text
                    
                    for field, value in data.items():
                        placeholder = f"{{{{{field}}}}}"
                        if placeholder in modified_text:
                            modified_text = modified_text.replace(placeholder, str(value))
                    
                    if modified_text != original_text:
                        t_elem.text = modified_text
                        
        except Exception:
            pass
    
    except Exception as e:
        st.warning(f"Advanced text box filling failed: {e}. Basic filling completed.")
    
    # Fill headers and footers
    try:
        for section in doc.sections:
            # Fill headers
            if section.header:
                for paragraph in section.header.paragraphs:
                    for field, value in data.items():
                        placeholder = f"{{{{{field}}}}}"
                        replace_text_in_paragraph(paragraph, placeholder, str(value))
            
            # Fill footers
            if section.footer:
                for paragraph in section.footer.paragraphs:
                    for field, value in data.items():
                        placeholder = f"{{{{{field}}}}}"
                        replace_text_in_paragraph(paragraph, placeholder, str(value))
    
    except Exception as e:
        st.error(f"Error filling headers/footers: {e}")
    
    return doc

# --- Main Application Logic ---

def main():
    st.warning('**DO NOT ENTER CUI OR PII INTO THIS SYSTEM - FOR BETA TESTING AND NON-OFFICIAL USE ONLY**')
    try:
        st.image("banner.png", use_container_width=True)
    except Exception as e:
        st.info("Info: `banner.png` not found. Skipping image banner.")

    st.markdown('<div class="main-header">📊 Document AI Field Filler</div>', unsafe_allow_html=True)
    st.markdown("**Transform your templates with AI-powered data filling! This tool supports PowerPoint, Word, and PDF documents. It will take unformatted data and conduct research, formatting, organization, data extraction, and place it in a pre-made template or bring your own!**")

    if 'fields' not in st.session_state:
        st.session_state.fields = []
    if 'field_locations' not in st.session_state:
        st.session_state.field_locations = []
    if 'ai_prompt' not in st.session_state:
        st.session_state.ai_prompt = ""

    st.markdown('<div class="step-container">', unsafe_allow_html=True)
    st.markdown("### 📁 Step 1: Choose Your Template")

    try:
        template_files = glob.glob("templates/*.*")
        template_options = ["Upload my own template"] + template_files
    except Exception as e:
        st.error(f"Could not scan templates directory: {e}")
        template_options = ["Upload my own template"]

    selected_template = st.selectbox("Select a template or upload your own:", options=template_options)
    source_file = None 
    template_name = None  # Track template name for prompt selection

    if selected_template == "Upload my own template":
        source_file = st.file_uploader("Choose your template file", type=['pptx', 'docx', 'pdf'])
        if source_file:
            template_name = source_file.name
    else:
        source_file = selected_template
        template_name = os.path.basename(selected_template)
    
    # Show prompt configuration info
    if template_name:
        prompt_config = load_prompt_config()
        if template_name in prompt_config.get("template_prompts", {}):
            st.success(f"✅ Custom prompt configured for: {template_name}")
        else:
            st.info(f"ℹ️ Using default prompt for: {template_name}")
    
    st.markdown('</div>', unsafe_allow_html=True)
    
    if source_file is not None:
        filename = source_file.name if hasattr(source_file, 'name') else source_file
        file_extension = filename.split('.')[-1].lower()

        with st.spinner('🔍 Analyzing template fields...'):
            if file_extension == 'pptx':
                st.session_state.fields, st.session_state.field_locations = analyze_powerpoint_fields(source_file)
            elif file_extension == 'docx':
                st.session_state.fields, st.session_state.field_locations = analyze_word_fields(source_file)
            elif file_extension == 'pdf':
                st.session_state.fields, st.session_state.field_locations = analyze_pdf_fields(source_file)
            else:
                st.error("Unsupported file type. Supported formats: PowerPoint (.pptx), Word (.docx), PDF (.pdf)")
                return

        if st.session_state.fields:
            st.markdown('<div class="success-box">', unsafe_allow_html=True)
            st.success(f"Found {len(st.session_state.fields)} placeholders in '{filename}'!")
            
            with st.expander("Click to see found fields and their locations"):
                col1, col2 = st.columns(2)
                with col1:
                    st.write("**Found Fields:**")
                    st.write(st.session_state.fields)
                with col2:
                    st.write("**Field Locations:**")
                    if st.session_state.field_locations:
                        for location in st.session_state.field_locations:
                            if file_extension == 'pdf':
                                st.write(f"• **{location['field']}** (Page {location['page']}, Type: {location['type']})")
                            elif file_extension == 'pptx':
                                st.write(f"• **{location['field']}** (Slide {location['slide']})")
                    else:
                        st.write("Location data not available for this document type.")

            st.markdown('</div>', unsafe_allow_html=True)

            # Create tabs for AI Generation and Manual Entry
            tab1, tab2 = st.tabs(["🤖 AI Generation", "✏️ Manual Entry"])
            
            # AI Generation Tab
            with tab1:
                st.markdown('<div class="step-container">', unsafe_allow_html=True)
                st.markdown("### 📝 Step 2: Enter Your Applicable Data Or Text. This can be formatted in any way, stream of thought, lists, sentences, etc. The more you provide the better your result will be. Any field on your template that is not covered will be TBD")
                project_data = st.text_area("Enter your data here:", height=200)
                
                # Image upload temporarily disabled for all formats
                uploaded_image = None

                st.markdown('</div>', unsafe_allow_html=True)

                if project_data.strip():
                    if st.button("🤖 Generate AI Prompt", type="primary", key="ai_prompt_btn"):
                        st.session_state.ai_prompt = generate_ai_prompt(st.session_state.fields, project_data, template_name)
                    
                    if st.session_state.ai_prompt:
                        st.markdown('<div class="step-container">', unsafe_allow_html=True)
                        st.markdown("### 📋 Step 3: Copy Prompt to AI")
                        st.info("Copy this prompt and paste it into your preferred AI assistant.")
                        
                        with st.expander("📄 Click to view the generated AI Prompt", expanded=False):
                            st.code(st.session_state.ai_prompt, language="text")
                        
                        # Use the original working copy component with popup
                        copy_component("📋 Copy Prompt to Clipboard", st.session_state.ai_prompt)
                        
                        # Show a temporary success message
                        if 'show_copy_success' not in st.session_state:
                            st.session_state.show_copy_success = False
                        
                        # Simple way to show feedback without breaking the copy function
                        st.info("💡 Click the button above to copy the prompt, then paste it into your AI assistant.")

                        st.markdown("**Quick Link to AI Service:**")
                        st.markdown(f'<a href="https://niprgpt.mil/" target="_blank" class="ai-button nipr-btn">🚀 Open NiprGPT</a>', unsafe_allow_html=True)
                        st.markdown('</div>', unsafe_allow_html=True)

                        st.markdown('<div class="step-container">', unsafe_allow_html=True)
                        st.markdown("### 🔄 Step 4: Paste AI Response & Generate")
                        ai_response = st.text_area("Paste the AI's JSON response here:", height=150)

                        if ai_response.strip():
                            try:
                                json_str_match = re.search(r'\{.*\}', ai_response, re.DOTALL)
                                if not json_str_match:
                                    raise json.JSONDecodeError("No JSON object found", "", 0)
                                json_data = json.loads(json_str_match.group(0))
                                st.success("✅ Valid JSON detected!")

                                if st.button("🚀 Generate Filled Document", type="primary", key="ai_generate_btn"):
                                    progress_container = st.container()
                                    with st.spinner('🔄 Filling template...'):
                                        if hasattr(source_file, 'seek'):
                                            source_file.seek(0)
                                        
                                        output_buffer = io.BytesIO()
                                        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
                                        
                                        if file_extension == 'pptx':
                                            prs = Presentation(source_file)
                                            filled_doc, _ = fill_powerpoint_with_data(prs, json_data, uploaded_image, progress_container)
                                            filled_doc.save(output_buffer)
                                            download_filename = f"filled_presentation_{timestamp}.pptx"
                                            mime_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
                                        elif file_extension == 'docx':
                                            filled_doc = fill_word_with_data(source_file, json_data)
                                            filled_doc.save(output_buffer)
                                            download_filename = f"filled_document_{timestamp}.docx"
                                            mime_type = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
                                        elif file_extension == 'pdf':
                                            filled_pdf_bytes, replacements = fill_pdf_with_data(source_file, json_data)
                                            if filled_pdf_bytes:
                                                output_buffer.write(filled_pdf_bytes)
                                                download_filename = f"filled_document_{timestamp}.pdf"
                                                mime_type = "application/pdf"
                                                progress_container.success(f"✅ PDF generated successfully! Made {replacements} replacements.")
                                            else:
                                                st.error("Failed to generate filled PDF")
                                                st.stop()
                                        
                                        if file_extension != 'pdf':
                                            progress_container.success("✅ Document generated successfully!")
                                        
                                        st.download_button(
                                            label=f"📥 Download Filled {file_extension.upper()}",
                                            data=output_buffer.getvalue(),
                                            file_name=download_filename,
                                            mime=mime_type
                                        )
                                        st.balloons()
                            except json.JSONDecodeError as e:
                                st.error(f"❌ Invalid JSON format: {e}")
                        st.markdown('</div>', unsafe_allow_html=True)
            
            # Manual Entry Tab
            with tab2:
                st.markdown('<div class="step-container">', unsafe_allow_html=True)
                st.markdown("### ✏️ Manual Entry - Fill Fields Directly")
                st.info(f"Found {len(st.session_state.fields)} fields to fill. Leave any field blank to keep the placeholder in the document.")
                
                # Initialize session state for manual entry if not exists
                if 'manual_entry_data' not in st.session_state:
                    st.session_state.manual_entry_data = {}
                
                # Create input fields for each found field
                st.markdown("**Fill in the fields below:**")
                
                # Split fields into two columns
                col1, col2 = st.columns(2)
                fields_per_column = (len(st.session_state.fields) + 1) // 2
                
                with col1:
                    for field in st.session_state.fields[:fields_per_column]:
                        # Use session state to preserve values
                        if field not in st.session_state.manual_entry_data:
                            st.session_state.manual_entry_data[field] = ""
                        
                        st.session_state.manual_entry_data[field] = st.text_input(
                            f"**{field}**",
                            value=st.session_state.manual_entry_data[field],
                            key=f"manual_field_{field}_1",
                            help=f"Enter value for {{{{ {field} }}}}"
                        )
                
                with col2:
                    for field in st.session_state.fields[fields_per_column:]:
                        # Use session state to preserve values
                        if field not in st.session_state.manual_entry_data:
                            st.session_state.manual_entry_data[field] = ""
                        
                        st.session_state.manual_entry_data[field] = st.text_input(
                            f"**{field}**",
                            value=st.session_state.manual_entry_data[field],
                            key=f"manual_field_{field}_2",
                            help=f"Enter value for {{{{ {field} }}}}"
                        )
                
                # Add utility buttons and generation
                st.markdown("---")
                col_clear, col_preview, col_generate = st.columns(3)
                
                with col_clear:
                    if st.button("🗑️ Clear All Fields", help="Clear all entered data"):
                        for field in st.session_state.fields:
                            st.session_state.manual_entry_data[field] = ""
                        st.rerun()
                
                with col_preview:
                    # Show preview of filled vs empty fields
                    filled_count = sum(1 for field in st.session_state.fields 
                                     if st.session_state.manual_entry_data.get(field, "").strip())
                    st.metric("Fields to Fill", f"{filled_count}/{len(st.session_state.fields)}")
                
                with col_generate:
                    if st.button("🚀 Generate Document", type="primary", key="manual_generate_btn"):
                        # Prepare data dictionary, excluding empty fields
                        manual_data = {}
                        filled_count = 0
                        
                        for field in st.session_state.fields:
                            value = st.session_state.manual_entry_data.get(field, "").strip()
                            if value:  # Only include non-empty fields
                                manual_data[field] = value
                                filled_count += 1
                        
                        st.info(f"Filling {filled_count} out of {len(st.session_state.fields)} fields. Empty fields will remain as placeholders.")
                        
                        progress_container = st.container()
                        with st.spinner('🔄 Generating document with manual entry...'):
                            if hasattr(source_file, 'seek'):
                                source_file.seek(0)
                            
                            output_buffer = io.BytesIO()
                            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
                            
                            if file_extension == 'pptx':
                                prs = Presentation(source_file)
                                filled_doc, _ = fill_powerpoint_with_data(prs, manual_data, None, progress_container)
                                filled_doc.save(output_buffer)
                                download_filename = f"manual_filled_presentation_{timestamp}.pptx"
                                mime_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
                            elif file_extension == 'docx':
                                filled_doc = fill_word_with_data(source_file, manual_data)
                                filled_doc.save(output_buffer)
                                download_filename = f"manual_filled_document_{timestamp}.docx"
                                mime_type = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
                            elif file_extension == 'pdf':
                                filled_pdf_bytes, replacements = fill_pdf_with_data(source_file, manual_data)
                                if filled_pdf_bytes:
                                    output_buffer.write(filled_pdf_bytes)
                                    download_filename = f"manual_filled_document_{timestamp}.pdf"
                                    mime_type = "application/pdf"
                                    progress_container.success(f"✅ PDF generated successfully! Made {replacements} replacements.")
                                else:
                                    st.error("Failed to generate filled PDF")
                                    st.stop()
                            
                            if file_extension != 'pdf':
                                progress_container.success("✅ Document generated successfully with manual entry!")
                            
                            st.download_button(
                                label=f"📥 Download Manual Filled {file_extension.upper()}",
                                data=output_buffer.getvalue(),
                                file_name=download_filename,
                                mime=mime_type
                            )
                            st.balloons()
                
                # Show a detailed preview of what will be filled
                with st.expander("📋 Preview of Field Mappings", expanded=False):
                    preview_data = []
                    for field in sorted(st.session_state.fields):
                        value = st.session_state.manual_entry_data.get(field, "").strip()
                        status = "✅ Will be filled" if value else "⚪ Will remain as placeholder"
                        preview_data.append({
                            "Field": f"{{{{{field}}}}}",
                            "Value": value if value else "(empty)",
                            "Status": status
                        })
                    
                    if preview_data:
                        preview_df = pd.DataFrame(preview_data)
                        st.dataframe(preview_df, use_container_width=True, hide_index=True)
                
                st.markdown('</div>', unsafe_allow_html=True)

        elif source_file is not None:
            st.markdown('<div class="warning-box">', unsafe_allow_html=True)
            st.warning("⚠️ No {{field_name}} placeholders found in your template!")
            
            if file_extension == 'pdf':
                st.info("""
                **PDF Tips:**
                - For text-based PDFs: Add placeholders like {{field_name}} in the text
                - For form-based PDFs: Use form field names that match your data fields
                - Some complex PDF structures may not be fully supported
                """)
            
            st.markdown('</div>', unsafe_allow_html=True)

    st.markdown("---")
    st.markdown("""
    <div style="text-align: center; color: #666; padding: 20px;">
        <p>🚀 Built for NIPR environments • No local installation required • Works in any browser</p>
        <p>📄 Supports PowerPoint (.pptx), Word (.docx), and PDF (.pdf) templates</p>
    </div>
    """, unsafe_allow_html=True)

if __name__ == "__main__":
    main()

